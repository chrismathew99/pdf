from flask import Flask, request, render_template, send_file
import os
from PyPDF2 import PdfFileReader, PdfFileWriter, PdfFileMerger
from pdf2image import convert_from_path
from docx import Document
from xlsxwriter import Workbook
from pptx import Presentation
import io
import zipfile
from werkzeug.utils import secure_filename
import logging

app = Flask(__name__)

# Ensure a folder for uploads exists
UPLOAD_FOLDER = 'uploads'
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# Set up logging
logging.basicConfig(level=logging.ERROR)
logger = logging.getLogger(__name__)

def allowed_file(filename, extensions=['pdf']):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in extensions


@app.route('/')
def home():
    return render_template('upload_page.html')


@app.route('/upload_pdf', methods=['POST'])
def upload_pdf():
    if request.method == 'POST':
        file = request.files['file']
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            try:
                # PDF upload and saving logic within try-except block
                file.save(os.path.join(app.config['UPLOAD_FOLDER'], filename))
                return render_template('manipulation_options_page.html', filename=filename)
            except Exception as e:
                # Log the exception and return a user-friendly error message
                logger.error(f"PDF processing failed: {str(e)}")
                return "Error processing PDF. Please ensure the file is not corrupted and try again."
    return 'File upload failed, ensure the file is a PDF.'


@app.route('/merge_pdfs', methods=['POST'])
def merge_pdfs():
    files = request.form.getlist('files')
    merger = PdfFileMerger()
    for pdf in files:
        full_path = os.path.join(app.config['UPLOAD_FOLDER'], pdf)
        if os.path.isfile(full_path):
            merger.append(full_path)
        else:
            return 'One or more files were not found.', 404
    output = io.BytesIO()
    merger.write(output)
    output.seek(0)
    return send_file(output, as_attachment=True, attachment_filename='merged.pdf')


@app.route('/split_pdf', methods=['POST'])
def split_pdf():
    filename = request.form['filename']
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    if not os.path.isfile(filepath):
        return 'File not found.', 404
    pdf = PdfFileReader(filepath)
    zipfile_obj = zipfile.ZipFile(io.BytesIO(), 'a', zipfile.ZIP_DEFLATED, False)
    for page in range(pdf.getNumPages()):
        pdf_writer = PdfFileWriter()
        pdf_writer.addPage(pdf.getPage(page))
        output = io.BytesIO()
        pdf_writer.write(output)
        zip_filename = f'split_{page}.pdf'
        zipfile_obj.writestr(zip_filename, output.getvalue())
    zipfile_io = zipfile_obj.fp
    zipfile_io.seek(0)
    zipfile_obj.close()
    return send_file(zipfile_io, mimetype='application/zip', as_attachment=True, attachment_filename='split_pdf.zip')


@app.route('/extract_text', methods=['POST'])
def extract_text():
    filename = request.form['filename']
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    if not os.path.isfile(filepath):
        return 'File not found.', 404
    pdf = PdfFileReader(filepath)
    text = ''
    for page in range(pdf.getNumPages()):
        text += pdf.getPage(page).extractText()
    return render_template('result_page.html', text=text)


@app.route('/convert_to_image', methods=['POST'])
def convert_to_image():
    filename = request.form['filename']
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    if not os.path.isfile(filepath):
        return 'File not found.', 404
    images = convert_from_path(filepath)
    image_bytes = io.BytesIO()
    images[0].save(image_bytes, format='JPEG')
    image_bytes.seek(0)
    return send_file(image_bytes, as_attachment=True, attachment_filename='image.jpg')


@app.route('/convert_to_excel', methods=['POST'])
def convert_to_excel():
    filename = request.form['filename']
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    if not os.path.isfile(filepath):
        return 'File not found.', 404
    workbook = Workbook(io.BytesIO())
    worksheet = workbook.add_worksheet()
    worksheet.write('A1', 'Extracted data from PDF')  # Placeholder, needs implementation
    workbook.close()
    output = workbook.filename
    output.seek(0)
    return send_file(output, as_attachment=True, attachment_filename=filename.replace('.pdf', '.xlsx'))


@app.route('/convert_to_word', methods=['POST'])
def convert_to_word():
    filename = request.form['filename']
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    if not os.path.isfile(filepath):
        return 'File not found.', 404
    doc = Document()
    doc.add_paragraph('Extracted data from PDF')  # Placeholder, needs implementation
    output = io.BytesIO()
    doc.save(output)
    output.seek(0)
    return send_file(output, as_attachment=True, attachment_filename=filename.replace('.pdf', '.docx'))


@app.route('/convert_to_powerpoint', methods=['POST'])
def convert_to_powerpoint():
    filename = request.form['filename']
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    if not os.path.isfile(filepath):
        return 'File not found.', 404
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Extracted data from PDF"  # Placeholder, needs implementation
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return send_file(output, as_attachment=True, attachment_filename=filename.replace('.pdf', '.pptx'))

if __name__ == '__main__':